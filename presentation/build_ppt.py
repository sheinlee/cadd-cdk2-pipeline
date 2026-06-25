#!/usr/bin/env python
"""Build a bilingual (中文/English) slide deck for the CDK2 CADD project (python-pptx)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

REPO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
RES = os.path.join(REPO, "results")
OUT = os.environ.get("PPT_OUT", os.path.join(REPO, "presentation", "CDK2_CADD_project_zh_en.pptx"))

NAVY = RGBColor(0x1F, 0x38, 0x64)
BLUE = RGBColor(0x2E, 0x75, 0xB6)
GREY = RGBColor(0x40, 0x40, 0x40)
LGREY = RGBColor(0x70, 0x70, 0x70)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Microsoft YaHei"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def _set(run, size, bold=False, color=GREY, font=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", font)


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def para(tf, text, size=18, bold=False, color=GREY, align=PP_ALIGN.LEFT, bullet=False,
         space_after=6, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    _set(r, size, bold, color)
    if bullet:
        p.level = 0
    return p


def band(slide, title_en, title_zh):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.25))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.45); tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title_en; _set(r, 26, True, WHITE)
    p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = title_zh; _set(r2, 18, False, RGBColor(0xCF, 0xE0, 0xF5))


def bullets_slide(title_en, title_zh, items):
    s = prs.slides.add_slide(BLANK)
    band(s, title_en, title_zh)
    tf = textbox(s, 0.7, 1.55, 12.0, 5.6)
    for i, (zh, en) in enumerate(items):
        para(tf, "•  " + zh, size=19, bold=True, color=NAVY, space_after=2, first=(i == 0))
        para(tf, "    " + en, size=14, color=LGREY, space_after=12)
    return s


def figure_slide(title_en, title_zh, takeaways, images):
    """images: list of (path, left_in, top_in, width_in)."""
    s = prs.slides.add_slide(BLANK)
    band(s, title_en, title_zh)
    tf = textbox(s, 0.7, 1.4, 12.0, 1.5)
    for i, (zh, en) in enumerate(takeaways):
        para(tf, "•  " + zh + "  —  " + en, size=14.5, bold=False, color=GREY, space_after=4,
             first=(i == 0))
    for path, l, t, w in images:
        if os.path.exists(path):
            s.shapes.add_picture(path, Inches(l), Inches(t), width=Inches(w))
    return s


# ---------- Slide 1: title ----------
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
tf = textbox(s, 0.9, 1.7, 11.5, 4.2)
para(tf, "CDK2 Structure-Based Virtual Screening", size=38, bold=True, color=WHITE, first=True)
para(tf, "Docking → MD → MM-GBSA → FEP   ·   an end-to-end, reproducible CADD pipeline",
     size=18, color=RGBColor(0xCF, 0xE0, 0xF5), space_after=18)
para(tf, "CDK2 基于结构的虚拟筛选 —— 对接→分子动力学→MM-GBSA→自由能微扰 全流程",
     size=20, bold=True, color=WHITE, space_after=4)
para(tf, "自主实践项目 / Self-initiated project    ·    李仕成 Shicheng Li", size=16,
     color=RGBColor(0xCF, 0xE0, 0xF5), space_after=4)
para(tf, "github.com/sheinlee/cadd-cdk2-pipeline", size=16, bold=True,
     color=RGBColor(0x9D, 0xC3, 0xE6))

# ---------- Slide 2: motivation ----------
bullets_slide("Motivation", "项目动机", [
    ("跨领域迁移:把已有的 OpenMM/MD 与生成式 ML 长板瞄准药物结合问题",
     "Transfer existing MD + generative-ML strengths to a drug-binding problem"),
    ("目标:独立搭建并跑通工业标准的 CADD 全流程,并对结果有批判性解读",
     "Build & run a full, reproducible, industry-standard CADD workflow end-to-end"),
    ("诚实定位:能力展示,而非新药发现的主张",
     "Honest framing: a capability demonstration, NOT a drug-discovery claim"),
    ("一切开源、脚本化、可一键复现",
     "Everything open-source, scripted, and reproducible from the repo"),
])

# ---------- Slide 3: target ----------
bullets_slide("Target — CDK2", "靶点 —— 细胞周期蛋白依赖性激酶 2", [
    ("CDK2:经典抗肿瘤激酶靶点;晶体结构 PDB 2R3I(1.28 Å,单体)",
     "Classic oncology kinase; crystal structure 2R3I (1.28 Å, monomeric)"),
    ("选它的理由:ATP 口袋深而规整(docking 友好)、抑制剂非共价(匹配方法)、公开数据多、体系小(单卡可跑)",
     "Deep well-defined ATP pocket, non-covalent inhibitors, rich public data, small system"),
    ("受体准备:还原氧化半胱氨酸 CSD177→CYS、重建无序 β3/αC loop、依共晶配体定义对接盒子",
     "Receptor prep: CSD177→CYS, rebuilt β3/αC loop, ATP-site box from the co-crystal ligand"),
])

# ---------- Slide 4: pipeline ----------
bullets_slide("Pipeline overview", "总体流程", [
    ("0 受体准备 (PDBFixer)  →  1 配体库:DUD-E 60 actives + 240 属性匹配 decoys",
     "Receptor prep → ligand library (DUD-E actives + matched decoys)"),
    ("2 对接 AutoDock Vina  →  富集 / ROC 验证",
     "Docking (AutoDock Vina) → enrichment / ROC validation"),
    ("3 短 MD (OpenMM, 4 ns, 显式溶剂)  →  pose 稳定性过滤",
     "Short explicit-solvent MD (OpenMM) → pose-stability filter"),
    ("4 MM-GBSA 端点重排 (AmberTools) + 每残基分解",
     "MM-GBSA endpoint rescoring (AmberTools) + per-residue decomposition"),
    ("验证track:相对结合 FEP (OpenFE) 对照实验",
     "Validation track: relative binding FEP (OpenFE) vs experiment"),
])

# ---------- Slide 5: docking enrichment ----------
figure_slide("Stage 1 — Docking enrichment", "第一步:对接富集验证", [
    ("300 配体(60 actives + 240 匹配 decoys)", "300 ligands, 60 actives / 240 matched decoys"),
    ("ROC-AUC 0.61,EF1% 3.3;原生 SCF 回docking −8.98(阳性对照)",
     "ROC-AUC 0.61, EF1% 3.3; native SCF redock −8.98 (control)"),
    ("结论:单构象刚性对接富集温和但真实 → top hits 需物理打分重排",
     "Modest-but-real enrichment → top hits motivate physics-based rescoring"),
], [
    (os.path.join(RES, "roc_curve.png"), 1.6, 3.1, 4.3),
    (os.path.join(RES, "score_distribution.png"), 7.0, 3.1, 4.8),
])

# ---------- Slide 6: MD + MM-GBSA ----------
figure_slide("Stage 2 — MD pose filter + MM-GBSA rescoring", "第二步:MD 稳定性过滤 + MM-GBSA 重排", [
    ("top 8 hits 各 4 ns 显式溶剂 MD:5/8 稳定,3 个漂移",
     "8 hits × 4 ns MD: 5/8 poses stable, 3 drift out"),
    ("MM-GBSA 把 Vina 的窄分数拉开 ~16 kcal/mol;best = 真实 active CHEMBL148580 (−53 kcal/mol)",
     "MM-GBSA spreads the scores; best hit is a true active (−53 kcal/mol)"),
    ("关键:MM-GBSA 排第2的 decoy 在 MD 中漂出口袋 → 被动力学揪出的假阳性",
     "A high-MM-GBSA decoy drifts out of the pocket — a false positive caught by dynamics"),
], [
    (os.path.join(RES, "docking_vs_mmgbsa.png"), 1.6, 3.05, 4.6),
    (os.path.join(RES, "rmsd_stability.png"), 6.9, 3.05, 5.1),
])

# ---------- Slide 7: hotspots ----------
figure_slide("Binding hotspots — per-residue MM-GBSA", "结合热点 —— 每残基能量分解", [
    ("自动识别配体周围口袋残基并做能量分解",
     "Pocket residues auto-detected; per-residue MM-GBSA decomposition"),
    ("定位到 hinge(Gln131/Asn132/Leu134)、催化 Lys33、Gly-rich 区 → 与已知 CDK2 药效团一致",
     "Localises to the hinge, catalytic Lys33 and Gly-rich loop — consistent with known CDK2 pharmacophores"),
], [
    (os.path.join(RES, "mmgbsa_decomposition.png"), 4.4, 2.95, 4.6),
])

# ---------- Slide 8: FEP ----------
figure_slide("Stage 3 — Relative binding FEP (the differentiator)", "第三步:相对结合 FEP(差异化王牌)", [
    ("OpenFE:OpenMM hybrid-topology RBFE + Hamiltonian replica exchange + MBAR",
     "OpenFE: OpenMM hybrid-topology RBFE + HREX + MBAR — the community-standard protocol"),
    ("JACS/Wang-2015 CDK2 基准,3 条 edge:MUE 0.71 / RMSE 1.1 / r 0.78;2/3 与实验差 ~0.2 kcal/mol",
     "JACS CDK2 benchmark, 3 edges: MUE 0.71 / RMSE 1.1 / r 0.78"),
    ("诚实:3 edge、repeats=1 为演示;lig_20→17 跑偏,如实报告、不挑结果",
     "Honest: 3-edge single-repeat demo; one edge misses — reported, not cherry-picked"),
], [
    (os.path.join(RES, "fep_correlation.png"), 4.5, 2.95, 4.4),
])

# ---------- Slide 9: engineering ----------
bullets_slide("Engineering — real problems solved", "工程能力 —— 真实踩坑与解决", [
    ("ChEMBL API 宕机/脆弱查询 → 改用稳定的 DUD-E 基准(保留 ChEMBL 备选)",
     "ChEMBL outage → switched to the static DUD-E benchmark (ChEMBL kept as fallback)"),
    ("OpenMM CUDA PTX 版本不匹配 → 把 cuda-version 钉到 12.6 对齐驱动",
     "OpenMM CUDA UNSUPPORTED_PTX_VERSION → pinned cuda-version=12.6 to the driver"),
    ("OpenBabel 从 pdbqt 错误感知键级(五价 N)→ 用 Meeko 直接重建分子",
     "OpenBabel mis-perceived bonds → reconstruct the ligand with Meeko"),
    ("MMPBSA print_res 语法不兼容 → 自动计算口袋残基 + 失败回退",
     "MMPBSA print_res incompatibility → auto-compute pocket residues with graceful fallback"),
    ("bash ( set -e ) || 陷阱(静默失效)→ 改用 && 链",
     "set -e silently disabled in a || context → replaced with explicit && chaining"),
])

# ---------- Slide 10: limitations ----------
bullets_slide("Honest assessment & limitations", "诚实评估与局限", [
    ("方法演示,非统计基准:库 300、FEP 3 edges、repeats=1、每 hit 单条 MD",
     "A methods demo, not a benchmark: small library, 3 FEP edges, single repeats"),
    ("MM-GBSA 绝对值未校准,仅相对排序有意义;FEP 那段已对照基准实验值",
     "MM-GBSA absolutes uncalibrated (relative only); FEP is validated against experiment"),
    ("生产级下一步:扩库、全 edge 网络 + ≥3 repeats、增强采样、刚性受体→柔性",
     "Production next steps: larger library, full edge network with ≥3 repeats, receptor flexibility"),
])

# ---------- Slide 11: summary ----------
s = prs.slides.add_slide(BLANK)
band(s, "Summary — skills demonstrated", "总结 —— 展示的能力")
tf = textbox(s, 0.7, 1.55, 12.0, 4.0)
para(tf, "• 全流程能力:结构准备 · 虚拟筛选 · 分子动力学 · 自由能(端点 MM-GBSA + 炼金术 FEP) · HPC/SLURM · Python 工程 · 可复现开源",
     size=17, bold=True, color=NAVY, first=True, space_after=4)
para(tf, "  End-to-end: structure prep · virtual screening · MD · free energy (MM-GBSA + alchemical FEP) · HPC/SLURM · reproducible open source",
     size=13, color=LGREY, space_after=14)
para(tf, "• 工具链 Toolchain:RDKit · AutoDock Vina · Meeko · OpenMM · AmberTools · OpenFE",
     size=17, bold=True, color=NAVY, space_after=14)
para(tf, "• 关键结果 Key results:docking AUC 0.61 / EF1% 3.3  ·  MM-GBSA 命中真实 active −53 kcal/mol  ·  FEP MUE 0.71 kcal/mol",
     size=17, bold=True, color=NAVY, space_after=14)
box = textbox(s, 0.7, 6.0, 12.0, 1.0)
para(box, "github.com/sheinlee/cadd-cdk2-pipeline", size=20, bold=True, color=BLUE, first=True)
para(box, "李仕成 Shicheng Li  ·  shichengli97@gmail.com", size=14, color=LGREY)

prs.save(OUT)
print("saved:", OUT, "slides:", len(prs.slides._sldIdLst))
